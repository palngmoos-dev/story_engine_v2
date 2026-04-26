import os
import re
from pptx import Presentation
from pptx.util import Pt

def produce_mirrored_slides():
    print("Starting Mirrored Ahn Eunhee Production...")
    
    root = r"d:\아름다운 여행 4.25"
    input_file = os.path.join(root, "Atelier", "input_space", "current_job.txt")
    template_path = os.path.join(root, "Atelier", "templates", "260516~0526 안은희 고객님 외 3인 고품격 프랑스 여행 9박 11일(신혜경,이두성 고객님 일정).pptx")
    output_path = r"C:\Users\moos\Desktop\AHN_EUNHEE_OFFICIAL_ITINERARY_V2.pptx"
    
    with open(input_file, "r", encoding="utf-8") as f:
        content = f.read()
    items = re.findall(r"(\d{2}\.\s\d{2}/\d{2}\s\w\s:.*)", content)
    
    prs = Presentation(template_path)
    
    # We want to use the same layout as slide 5 (index 5)
    # The layout object must come from the SAME presentation object we are adding to.
    target_layout = prs.slides[5].slide_layout
    
    for i, item in enumerate(items):
        match = re.match(r"(\d{2}\.\s\d{2}/\d{2}\s\w)\s:(.*)", item)
        if not match: continue
        
        date_part = match.group(1).strip()
        desc_part = match.group(2).strip()
        
        # Add slide based on the target layout
        slide = prs.slides.add_slide(target_layout)
        
        # Find placeholders or shapes to replace
        for shape in slide.shapes:
            if shape.is_placeholder:
                ph = shape.placeholder_format
                if ph.idx == 2: # Title
                    shape.text = f"❖ {date_part} 일정 안내"
                elif ph.idx == 3: # Page
                    shape.text = f"{i+1}/{len(items)}"
                elif ph.idx == 1: # Content
                    shape.text = desc_part.replace("/", "\n• ")
            elif "텍스트 개체 틀 1" in shape.name or shape.shape_type == 17:
                # If it's the main text box
                shape.text_frame.text = desc_part.replace("/", "\n• ")
                # Apply Master Guidelines
                for p in shape.text_frame.paragraphs:
                    p.font.size = Pt(24)
                    p.font.name = "나눔바른펜OTF"

        print(f"Mirrored slide for {date_part}")

    prs.save(output_path)
    print(f"Mirrored Production Complete: {output_path}")

if __name__ == "__main__":
    produce_mirrored_slides()
