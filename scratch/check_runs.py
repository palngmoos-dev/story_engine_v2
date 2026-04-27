import sys
from pptx import Presentation

# Force UTF-8 for output
sys.stdout.reconfigure(encoding='utf-8')

def analyze(path):
    prs = Presentation(path)
    slides = list(prs.slides)
    for i, slide in enumerate(slides[:5]):
        print(f"Slide {i+1}")
        for j, shape in enumerate(slide.shapes):
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    for r in p.runs:
                        print(f"  Shape {j+1} Run: [{r.text}]")
        print("-" * 20)

if __name__ == "__main__":
    analyze(r"D:\temp\template.pptx")
