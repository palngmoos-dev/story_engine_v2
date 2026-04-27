import sys
import os
import shutil
import re
from pptx import Presentation

# Ensure dependencies are available
sys.path.insert(0, r'C:\Users\moos\AppData\Roaming\Python\Python314\site-packages')

def smart_replace(p):
    """Correctly replace itinerary data in a paragraph, even if split into runs"""
    original_text = p.text
    text = original_text
    
    # 1. Names: Remove the extra two customers
    # Handle variations in spacing/commas
    text = re.sub(r'안은희\s*,\s*장혁수\s*,\s*신혜경\s*,\s*이두성', '안은희,장혁수', text)
    
    # 2. Date Range: 5/16 ~ 5/26 -> 5/16 ~ 6/7
    # Using flexible regex to catch various formats
    text = re.sub(r'5월\s*16일.*?~.*?5월\s*26일', '5월 16일(토요일) ~ 6월 7일(일요일)', text)
    
    # 3. Trip Duration: 9박 11일 -> 22박 23일
    text = text.replace('9박 11일', '22박 23일')
    
    # 4. If we missed the combined names, try individual removals
    if '신혜경' in text or '이두성' in text:
        text = text.replace('신혜경', '').replace('이두성', '')
        # Clean up stray commas
        text = text.replace(',,', ',').strip(' ,')
    
    if text != original_text:
        # Note: setting p.text might reset some font styles, 
        # but it's the most reliable way to fix split runs.
        p.text = text
        return True
    return False

def process_shapes(shapes):
    updated = 0
    for shape in shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                if smart_replace(paragraph):
                    updated += 1
        elif shape.shape_type == 19: # Table
            for row in shape.table.rows:
                for cell in row.cells:
                    for paragraph in cell.text_frame.paragraphs:
                        if smart_replace(paragraph):
                            updated += 1
        elif shape.shape_type == 6: # Group
            updated += process_shapes(shape.shapes)
    return updated

def main():
    template_path = r"D:\temp\template.pptx"
    output_path = r"D:\아름다운 여행 4.25\지금 작업 4.25\260516~0607 안은희 고객님 고품격 프랑스 여행 일정_최종수정본.pptx"
    
    print(f"Loading template from: {template_path}")
    if not os.path.exists(template_path):
        print("Template not found!")
        return

    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    shutil.copy2(template_path, output_path)
    
    prs = Presentation(output_path)
    print(f"Processing {len(prs.slides)} slides...")
    
    total_updated = 0
    for i, slide in enumerate(prs.slides):
        slide_updates = process_shapes(slide.shapes)
        total_updated += slide_updates
        if slide_updates > 0:
            print(f"  Slide {i+1}: Updated {slide_updates} items")
            
    prs.save(output_path)
    print(f"\nSUCCESS: Task completed.")
    print(f"Total items updated: {total_updated}")
    print(f"Final file: {output_path}")

if __name__ == "__main__":
    main()
