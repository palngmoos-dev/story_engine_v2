import sys
sys.path.insert(0, r'C:\Users\moos\AppData\Roaming\Python\Python314\site-packages')

from pptx import Presentation
import shutil
import os

def find_and_replace_text(prs, replacements):
    """Find and replace text throughout the presentation correctly"""
    total_replacements = 0

    for slide_idx, slide in enumerate(prs.slides):
        slide_replacements = 0

        for shape in slide.shapes:
            # 1. Text Frames
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        original_text = run.text
                        new_text = original_text
                        
                        # Correct loop: apply each replacement to the CURRENT text
                        for old_key, new_val in replacements.items():
                            if old_key and old_key in new_text:
                                new_text = new_text.replace(old_key, new_val)
                        
                        if new_text != original_text:
                            run.text = new_text
                            slide_replacements += 1
                            total_replacements += 1

            # 2. Tables
            elif shape.shape_type == 19:  # Table
                for row in shape.table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                original_text = run.text
                                new_text = original_text
                                for old_key, new_val in replacements.items():
                                    if old_key and old_key in new_text:
                                        new_text = new_text.replace(old_key, new_val)
                                if new_text != original_text:
                                    run.text = new_text
                                    slide_replacements += 1
                                    total_replacements += 1

            # 3. Grouped shapes (recursive check)
            elif shape.shape_type == 6:  # Group
                total_replacements += process_group_shape(shape, replacements)

    return total_replacements

def process_group_shape(group_shape, replacements):
    count = 0
    for shape in group_shape.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    original_text = run.text
                    new_text = original_text
                    for old_key, new_val in replacements.items():
                        if old_key and old_key in new_text:
                            new_text = new_text.replace(old_key, new_val)
                    if new_text != original_text:
                        run.text = new_text
                        count += 1
        elif shape.shape_type == 6:
            count += process_group_shape(shape, replacements)
    return count

def main():
    # Define paths
    template_path = r"D:\temp\template.pptx"
    output_path = r"D:\아름다운 여행 4.25\지금 작업 4.25\260516~0607 안은희 고객님 고품격 프랑스 여행 일정_업데이트.pptx"

    # Ensure output directory exists
    os.makedirs(os.path.dirname(output_path), exist_ok=True)

    # Copy template to output
    shutil.copy2(template_path, output_path)
    print(f"Template copied to: {output_path}")

    # Load presentation
    prs = Presentation(output_path)
    print(f"Loaded presentation with {len(prs.slides)} slides")

    # Define specific replacements
    replacements = {
        '안은희,장혁수,신혜경,이두성': '안은희,장혁수',
        '5월 16일(토요일) ~ 5월 26일(화요일)': '5월 16일(토요일) ~ 6월 7일(일요일)',
        '총 9박 11일': '총 22박 23일',
        # Don't replace the footer with itself if it causes issues, but the logic is now fixed anyway.
        # Removing the footer from replacements since we don't need to change it.
    }

    print("Performing text replacements (CORRECTED)...")
    total_replacements = find_and_replace_text(prs, replacements)
    print(f"Total replacements made: {total_replacements}")

    # Save the presentation
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")

if __name__ == "__main__":
    main()
