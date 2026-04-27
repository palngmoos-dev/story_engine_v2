import sys
sys.path.insert(0, r'C:\Users\moos\AppData\Roaming\Python\Python314\site-packages')

from pptx import Presentation
import shutil
import os

def find_and_replace_text(prs, replacements):
    """Find and replace text throughout the presentation correctly"""
    total_replacements = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        original_text = run.text
                        new_text = original_text
                        for old_key, new_val in replacements.items():
                            if old_key and old_key in new_text:
                                new_text = new_text.replace(old_key, new_val)
                        if new_text != original_text:
                            run.text = new_text
                            total_replacements += 1
            elif shape.shape_type == 19:  # Table
                for row in shape.table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                original_text = run.text
                                new_text = original_text
                                for old_key, new_val in replacements.items():
                                    if old_key and old_key in new_text:
                                        new_text = new_text.replace(old_key, new_val)
                                if new_text != original_text:
                                    run.text = new_text
                                    total_replacements += 1
            elif shape.shape_type == 6:  # Group
                total_replacements += process_group_shape(shape, replacements)
    return total_replacements

def process_group_shape(group_shape, replacements):
    count = 0
    for shape in group_shape.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    original_text = run.text
                    new_text = original_text
                    for old_key, new_val in replacements.items():
                        if old_key and old_key in new_text:
                            new_text = new_text.replace(old_key, new_val)
                    if new_text != original_text:
                        run.text = new_text
                        count += 1
        elif shape.shape_type == 6:
            count += process_group_shape(shape, replacements)
    return count

def main():
    template_path = r"D:\temp\template.pptx"
    # Using a NEW filename to avoid permission errors
    output_path = r"D:\아름다운 여행 4.25\지금 작업 4.25\260516~0607 안은희 고객님 고품격 프랑스 여행 일정_수정본.pptx"

    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    shutil.copy2(template_path, output_path)
    print(f"Template copied to: {output_path}")

    prs = Presentation(output_path)
    replacements = {
        '안은희,장혁수,신혜경,이두성': '안은희,장혁수',
        '5월 16일(토요일) ~ 5월 26일(화요일)': '5월 16일(토요일) ~ 6월 7일(일요일)',
        '총 9박 11일': '총 22박 23일',
    }

    print("Performing text replacements (CORRECTED)...")
    total_replacements = find_and_replace_text(prs, replacements)
    print(f"Total replacements made: {total_replacements}")

    prs.save(output_path)
    print(f"Presentation saved successfully: {output_path}")

if __name__ == "__main__":
    main()
