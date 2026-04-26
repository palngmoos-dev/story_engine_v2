import os
from pptx import Presentation
import json

def analyze_slide_structure(pptx_path, slide_idx=5):
    print(f"Analyzing {pptx_path} Slide {slide_idx}...")
    prs = Presentation(pptx_path)
    if slide_idx >= len(prs.slides):
        print(f"Slide index {slide_idx} out of range (Total: {len(prs.slides)})")
        return
    
    slide = prs.slides[slide_idx]
    structure = []
    for shape in slide.shapes:
        info = {
            "name": shape.name,
            "type": str(shape.shape_type),
            "left": shape.left.pt,
            "top": shape.top.pt,
            "width": shape.width.pt,
            "height": shape.height.pt,
        }
        if shape.has_text_frame:
            info["text"] = shape.text_frame.text
            if shape.text_frame.paragraphs:
                p = shape.text_frame.paragraphs[0]
                if p.runs:
                    r = p.runs[0]
                    info["font"] = {
                        "name": r.font.name,
                        "size": r.font.size.pt if r.font.size else None,
                        "bold": r.font.bold
                    }
        structure.append(info)
    
    with open("template_structure.json", "w", encoding="utf-8") as f:
        json.dump(structure, f, indent=4, ensure_ascii=False)
    print("Structure saved to template_structure.json")

if __name__ == "__main__":
    tpl_path = r"d:\아름다운 여행 4.25\Atelier\templates\260516~0526 안은희 고객님 외 3인 고품격 프랑스 여행 9박 11일(신혜경,이두성 고객님 일정).pptx"
    analyze_slide_structure(tpl_path)
