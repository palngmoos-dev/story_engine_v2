import os
import re
import glob
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from PIL import Image

# ==========================================
# AHN EUNHEE MASTERPIECE FINAL PRODUCTION
# Following Master Guidelines (44pt/24pt)
# Including Smart Fit Image Engine
# ==========================================

class MasterpieceEngine:
    def __init__(self, template_path, asset_dir):
        self.prs = Presentation(template_path)
        self.asset_dir = asset_dir
        self.layout = self.prs.slides[5].slide_layout # Mirroring slide 5
        self.cropped_dir = os.path.join(asset_dir, "master_cropped")
        if not os.path.exists(self.cropped_dir): os.makedirs(self.cropped_dir)

    def _crop_and_get_ratio(self, photo_path):
        img = Image.open(photo_path)
        width, height = img.size
        # Smart crop: if it's a map (wide), crop the sidebar. If it's a photo, keep it.
        if "map" in photo_path.lower():
            left = int(width * 0.35)
            img = img.crop((left, 0, width, height))
        
        save_path = os.path.join(self.cropped_dir, os.path.basename(photo_path))
        img.save(save_path)
        return save_path, img.size[0] / img.size[1]

    def _add_image_smart_fit(self, slide, photo_path):
        try:
            path, ratio = self._crop_and_get_ratio(photo_path)
            # Define image area (centered below title)
            max_w = Inches(11.0)
            max_h = Inches(4.5)
            left_base = Inches(1.1)
            top_base = Inches(2.2)

            calc_w = max_h * ratio
            if calc_w > max_w:
                final_w = max_w
                final_h = max_w / ratio
            else:
                final_w = calc_w
                final_h = max_h

            offset_x = left_base + (max_w - final_w) / 2
            offset_y = top_base + (max_h - final_h) / 2
            
            slide.shapes.add_picture(path, offset_x, offset_y, width=final_w, height=final_h)
        except Exception as e:
            print(f"Error adding image {photo_path}: {e}")

    def create_itinerary_slide(self, date_part, desc_part, image_keyword=None):
        slide = self.prs.slides.add_slide(self.layout)
        
        # 1. Title (Pt 44)
        for shape in slide.shapes:
            if shape.is_placeholder and shape.placeholder_format.idx == 2:
                shape.text = f"❖ {date_part} 일정 안내"
                for p in shape.text_frame.paragraphs:
                    p.font.size = Pt(44)
                    p.font.name = "나눔바른펜OTF"
                    p.font.bold = True
                    p.font.color.rgb = RGBColor(68, 84, 106)
                break
        
        # 2. Content (Pt 24)
        # Find main text box or Placeholder 1
        for shape in slide.shapes:
            if (shape.is_placeholder and shape.placeholder_format.idx == 1) or ("텍스트 개체 틀 1" in shape.name):
                shape.text = desc_part.replace("/", "\n• ")
                for p in shape.text_frame.paragraphs:
                    p.font.size = Pt(24)
                    p.font.name = "나눔바른펜OTF"
                    p.space_after = Pt(10)
                break

        # 3. Add Image if keyword matches
        if image_keyword:
            matches = glob.glob(os.path.join(self.asset_dir, f"*{image_keyword}*photo*"))
            if not matches:
                matches = glob.glob(os.path.join(self.asset_dir, f"*{image_keyword}*"))
            
            if matches:
                self._add_image_smart_fit(slide, matches[0])

    def save(self, output_path):
        self.prs.save(output_path)
        print(f"Masterpiece Saved: {output_path}")

def produce_ultimate():
    root = r"d:\아름다운 여행 4.25"
    asset_dir = r"C:\Users\moos\.gemini\antigravity\brain\2fb7398c-27ff-4ae5-b026-1f183e98301a"
    template_path = os.path.join(root, "Atelier", "templates", "260516~0526 안은희 고객님 외 3인 고품격 프랑스 여행 9박 11일(신혜경,이두성 고객님 일정).pptx")
    input_file = os.path.join(root, "Atelier", "input_space", "current_job.txt")
    output_path = r"C:\Users\moos\Desktop\AHN_EUNHEE_MASTERPIECE_OFFICIAL.pptx"

    with open(input_file, "r", encoding="utf-8") as f:
        content = f.read()
    items = re.findall(r"(\d{2}\.\s\d{2}/\d{2}\s\w)\s:(.*)", content)

    # Keyword mapping for images
    keyword_map = {
        "05/17": "nice",
        "05/19": "st_paul_de_vence",
        "05/20": "eze",
        "05/21": "cassis",
        "05/24": "arles",
        "05/25": "gordes",
        "06/02": "milan",
        "06/04": "stresa"
    }

    engine = MasterpieceEngine(template_path, asset_dir)
    
    # Process only the specific 09-23 section for the masterpiece
    for item in items:
        date_tag = item[0][-5:] # e.g. 05/24
        keyword = keyword_map.get(date_tag)
        engine.create_itinerary_slide(item[0], item[1], keyword)
    
    engine.save(output_path)

if __name__ == "__main__":
    produce_ultimate()
