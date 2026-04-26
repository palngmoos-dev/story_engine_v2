import os
from pptx import Presentation
import json
import sys

# PPTX Forensic Scanner V2
# Fixed slicing issue and added more robust attribute extraction.

sys.stdout.reconfigure(encoding='utf-8')

def forensic_scan(pptx_path):
    print(f"Forensic Scanning: {pptx_path}")
    prs = Presentation(pptx_path)
    scan_results = {
        "file_name": os.path.basename(pptx_path),
        "slides": []
    }

    count = 0
    for slide in prs.slides:
        if count >= 10: break
        count += 1
        
        slide_data = {"index": count, "shapes": []}
        for shape in slide.shapes:
            shape_info = {
                "name": shape.name,
                "type": str(shape.shape_type),
                "left": shape.left, "top": shape.top, "width": shape.width, "height": shape.height,
                "rotation": shape.rotation
            }
            
            # Line/Border style
            if hasattr(shape, "line"):
                try:
                    shape_info["line"] = {
                        "color": str(shape.line.color.rgb) if hasattr(shape.line.color, 'rgb') else None,
                        "width": shape.line.width.pt if shape.line.width else None
                    }
                except: pass
                
            # Fill style
            if hasattr(shape, "fill"):
                try:
                    shape_info["fill"] = {
                        "transparency": shape.fill.transparency if hasattr(shape.fill, 'transparency') else 0
                    }
                except: pass

            # Text properties
            if shape.has_text_frame:
                text_props = []
                for p in shape.text_frame.paragraphs:
                    for r in p.runs:
                        text_props.append({
                            "font_name": r.font.name,
                            "font_size": r.font.size.pt if r.font.size else None,
                            "bold": r.font.bold,
                            "color": str(r.font.color.rgb) if hasattr(r.font.color, 'rgb') else None
                        })
                shape_info["text_properties"] = text_props
                
            slide_data["shapes"].append(shape_info)
        scan_results["slides"].append(slide_data)
        
    return scan_results

if __name__ == "__main__":
    root = r"d:\아름다운 여행 4.25"
    dirs = [d for d in os.listdir(root) if os.path.isdir(os.path.join(root, d))]
    target_dir = None
    for d in dirs:
        if any(c in d for c in ["ߺ", "Gyeom", "견본"]):
            target_dir = os.path.join(root, d)
            break
            
    if not target_dir:
        print("Template directory not found.")
        sys.exit(1)
        
    pptx_files = [f for f in os.listdir(target_dir) if f.endswith(".pptx")]
    master_pptx = os.path.join(target_dir, sorted(pptx_files, key=lambda x: os.path.getsize(os.path.join(target_dir, x)), reverse=True)[0])
    
    try:
        results = forensic_scan(master_pptx)
        with open("forensic_design_scan.json", "w", encoding="utf-8") as f:
            json.dump(results, f, ensure_ascii=False, indent=4)
        print("Forensic scan complete.")
    except Exception as e:
        print(f"Error during scan: {e}")
