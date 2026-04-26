import os
import re
from pptx import Presentation
from pptx.util import Pt

def produce_final_slides():
    print("Starting Ahn Eunhee Masterpiece Production...")
    
    # 1. Paths
    root = r"d:\아름다운 여행 4.25"
    input_file = os.path.join(root, "Atelier", "input_space", "current_job.txt")
    template_path = os.path.join(root, "Atelier", "templates", "260516~0526 안은희 고객님 외 3인 고품격 프랑스 여행 9박 11일(신혜경,이두성 고객님 일정).pptx")
    output_path = r"C:\Users\moos\Desktop\AHN_EUNHEE_OFFICIAL_ITINERARY.pptx"
    
    # 2. Extract Itinerary
    with open(input_file, "r", encoding="utf-8") as f:
        content = f.read()
    
    # Match items 09 to 23
    items = re.findall(r"(\d{2}\.\s\d{2}/\d{2}\s\w\s:.*)", content)
    
    if not items:
        print("No itinerary items found!")
        return

    # 3. Load Template and Layout
    prs = Presentation(template_path)
    # We will use slide 5 as a reference for formatting if needed, 
    # but for simplicity in V3, we'll just add slides with a clean layout.
    # To truly "mirror", we should copy shapes.
    
    # Let's try to clone slide 5's style for each item.
    ref_slide = prs.slides[5]
    
    # Create a new presentation to keep it clean
    new_prs = Presentation()
    new_prs.slide_width = prs.slide_width
    new_prs.slide_height = prs.slide_height
    
    for i, item in enumerate(items):
        # Parse item
        match = re.match(r"(\d{2}\.\s\d{2}/\d{2}\s\w)\s:(.*)", item)
        if not match: continue
        
        date_part = match.group(1).strip()
        desc_part = match.group(2).strip()
        
        # Add slide
        slide = new_prs.slides.add_slide(new_prs.slide_layouts[6]) # Blank layout
        
        # Add Title (Master Guidelines: Pt 44)
        title_box = slide.shapes.add_textbox(Pt(50), Pt(30), Pt(700), Pt(60))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"❖ {date_part} 일정 안내"
        p.font.name = "나눔바른펜OTF"
        p.font.size = Pt(44)
        p.font.bold = True
        
        # Add Content (Master Guidelines: Pt 24)
        content_box = slide.shapes.add_textbox(Pt(50), Pt(100), Pt(700), Pt(400))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        # Split content by '/' to create bullet points if possible
        points = desc_part.split("/")
        for point in points:
            p = tf.add_paragraph()
            p.text = f"• {point.strip()}"
            p.font.name = "나눔바른펜OTF"
            p.font.size = Pt(24)
            p.space_after = Pt(10)

        print(f"Added slide for {date_part}")

    # 4. Save
    new_prs.save(output_path)
    print(f"Production Complete: {output_path}")

if __name__ == "__main__":
    produce_final_slides()
