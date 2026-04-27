from pptx import Presentation
import sys

def analyze_pptx(file_path):
    """Analyze the structure of a PowerPoint file"""
    try:
        prs = Presentation(file_path)
        print(f"Analyzing: {file_path}")
        print(f"Number of slides: {len(prs.slides)}")
        print("=" * 50)

        for i, slide in enumerate(prs.slides):
            print(f"Slide {i+1}:")
            print(f"  - Number of shapes: {len(slide.shapes)}")

            for j, shape in enumerate(slide.shapes):
                if hasattr(shape, "text") and shape.text.strip():
                    print(f"    Shape {j+1}: {shape.text[:100]}{'...' if len(shape.text) > 100 else ''}")
                elif shape.shape_type == 13:  # Picture
                    print(f"    Shape {j+1}: Picture")
                elif shape.shape_type == 1:  # AutoShape
                    print(f"    Shape {j+1}: AutoShape")
                else:
                    print(f"    Shape {j+1}: Type {shape.shape_type}")
            print()

    except Exception as e:
        print(f"Error analyzing {file_path}: {e}")
        return False

    return True

if __name__ == "__main__":
    template_path = r"D:\아름다운 여행 4.25\견본 4.25\260516~0526 안은희 고객님 외 3인 고품격 프랑스 여행 9박 11일(신혜경,이두성 고객님 일정).pptx"
    analyze_pptx(template_path)