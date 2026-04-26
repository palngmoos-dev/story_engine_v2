import os
from pptx import Presentation
import json
import sys

sys.stdout.reconfigure(encoding='utf-8')

def analyze_pptx(file_name):
    prs = Presentation(file_name)
    analysis = {
        "file_name": file_name,
        "total_slides": len(prs.slides),
        "slides": []
    }

    for i, slide in enumerate(prs.slides):
        slide_info = {
            "slide_index": i + 1,
            "layout": slide.slide_layout.name,
            "shapes": []
        }
        
        for shape in slide.shapes:
            shape_info = {
                "type": str(shape.shape_type),
                "name": shape.name,
                "left": shape.left,
                "top": shape.top,
                "width": shape.width,
                "height": shape.height
            }
            
            if hasattr(shape, "text"):
                shape_info["text"] = shape.text
            
            slide_info["shapes"].append(shape_info)
            
        analysis["slides"].append(slide_info)
    
    return analysis

if __name__ == "__main__":
    pptx_path = "target.pptx"
    if not os.path.exists(pptx_path):
        print(f"File not found: {pptx_path}")
        sys.exit(1)
        
    try:
        data = analyze_pptx(pptx_path)
        with open("pptx_analysis.json", 'w', encoding='utf-8') as f:
            json.dump(data, f, ensure_ascii=False, indent=4)
        print(f"Analysis saved to pptx_analysis.json")
    except Exception as e:
        print(f"Error: {e}")
