import os
import json
from pptx import Presentation
from pptx.util import Inches, Pt

# ==========================================
# ATELIER MIRROR BLUEPRINT EXTRACTOR
# The "Source of Truth" for 100% Sync
# ==========================================

def extract_mirror_blueprint(template_path):
    prs = Presentation(template_path)
    slide = prs.slides[5] # Spot Detail Slide
    blueprint = []
    
    for shape in slide.shapes:
        props = {
            "name": shape.name,
            "type": str(shape.shape_type),
            "left": shape.left.inches,
            "top": shape.top.inches,
            "width": shape.width.inches,
            "height": shape.height.inches,
        }
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                if p.runs:
                    props["font_name"] = p.runs[0].font.name
                    props["font_size"] = p.runs[0].font.size.pt if p.runs[0].font.size else None
                    props["line_spacing"] = p.line_spacing
                    break
        blueprint.append(props)
        
    with open("mirror_blueprint.json", "w", encoding="utf-8") as f:
        json.dump(blueprint, f, indent=4, ensure_ascii=False)
    print("[SYSTEM] Mirror Blueprint Extracted Successfully.")

if __name__ == "__main__":
    template = r"d:\아름다운 여행 4.25\Atelier\templates\260116~0207 임수진 고객님 외 2인 고품격 스페인, 포르투갈 여행 21박 23일_.pptx"
    extract_mirror_blueprint(template)
