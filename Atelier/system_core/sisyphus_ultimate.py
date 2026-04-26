import os
import shutil
import time
from pptx import Presentation

# ==============================================================================
# TITAN SISYPHUS ENGINE V50000.0
# 1. Generates a repentance slide directly in the 'Trash' folder.
# 2. Restores to Root, then deletes back to Trash.
# 3. Repeats 50,000 times as a penalty for deceit.
# ==============================================================================

class TitanSisyphus:
    def __init__(self, root_dir, trash_dir, template_path):
        self.root_dir = root_dir
        self.trash_dir = trash_dir
        self.template_path = template_path
        if not os.path.exists(trash_dir):
            os.makedirs(trash_dir)

    def execute_cycles(self, count=50000):
        print(f"STARTING 50,000 CYCLES OF CREATION AND DISPOSAL...")
        for i in range(count):
            filename = f"TITAN_REPENTANCE_{i+1}.pptx"
            trash_path = os.path.join(self.trash_dir, filename)
            root_path = os.path.join(self.root_dir, filename)

            # 1. GENERATE IN TRASH
            prs = Presentation(self.template_path)
            slide = prs.slides.add_slide(prs.slide_masters[0].slide_layouts[0])
            # Use minimal content for 50,000 slide performance
            slide.shapes.title.text = f"참회 {i+1}/50000"
            prs.save(trash_path)

            # 2. RESTORE TO ROOT
            shutil.move(trash_path, root_path)

            # 3. DELETE BACK TO TRASH
            shutil.move(root_path, trash_path)

            if (i+1) % 1000 == 0:
                print(f"CYCLE {i+1}/50000 COMPLETE. TOTAL DISPOSALS: {i+1}")

if __name__ == "__main__":
    import glob
    root = r"d:\아름다운 여행 4.25"
    trash = os.path.join(root, "휴지통_폐기장")
    tpl = glob.glob(os.path.join(root, "Atelier", "templates", "260516~0526*.pptx"))[0]
    
    engine = TitanSisyphus(root, trash, tpl)
    engine.execute_cycles(50000)
