import os
import glob
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# ==============================================================================
# SMART FIT ENGINE V2.0
# 1. Crops the sidebar from maps.
# 2. Calculates the new aspect ratio to prevent distortion.
# 3. Fits the map perfectly onto the slide.
# ==============================================================================

def crop_and_get_ratio(photo_path, cropped_dir):
    if not os.path.exists(photo_path):
        return None, 1.0
    
    img = Image.open(photo_path)
    width, height = img.size
    # Crop 35% from left
    left = int(width * 0.35)
    cropped_img = img.crop((left, 0, width, height))
    
    out_name = os.path.basename(photo_path)
    save_path = os.path.join(cropped_dir, out_name)
    cropped_img.save(save_path)
    
    new_width, new_height = cropped_img.size
    return save_path, new_width / new_height

class SmartFitMapEngine:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.33)
        self.prs.slide_height = Inches(7.5)
        
    def _add_picture_smart_fit(self, slide, photo_path, ratio):
        # Max dimensions for the map
        max_w = Inches(12.3)
        max_h = Inches(5.8)
        left_m = Inches(0.5)
        top_m = Inches(0.8)

        # Calculate dimensions based on ratio
        # ratio = w / h -> w = h * ratio
        calc_w = max_h * ratio
        if calc_w > max_w:
            # Fit to width
            final_w = max_w
            final_h = max_w / ratio
        else:
            # Fit to height
            final_w = calc_w
            final_h = max_h

        # Center in the available space
        offset_x = left_m + (max_w - final_w) / 2
        offset_y = top_m + (max_h - final_h) / 2
        
        slide.shapes.add_picture(photo_path, offset_x, offset_y, width=final_w, height=final_h)

    def _create_slide(self, title, schedule, map_url, photo_pattern, asset_dir, cropped_dir):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(12), Inches(0.6))
        tf = title_box.text_frame
        run = tf.paragraphs[0].add_run()
        run.text = f"❖ {title}"
        run.font.name = 'HY동녘M'
        run.font.size = Pt(44)
        run.font.color.rgb = RGBColor(68, 84, 106)

        # Process Image
        assets = glob.glob(os.path.join(asset_dir, photo_pattern))
        if assets:
            save_path, ratio = crop_and_get_ratio(assets[0], cropped_dir)
            if save_path:
                self._add_picture_smart_fit(slide, save_path, ratio)
                # Hyperlink still on the title or invisible overlay? Let's add it to the image if possible.
                # Actually, the picture itself
                for shape in slide.shapes:
                    if shape.shape_type == 13:
                        shape.click_action.hyperlink.address = map_url

        # Minimal Text
        text_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.7))
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = schedule.replace("\n", " | ")
        run.font.name = '휴먼매직체'
        run.font.size = Pt(24)

    def produce(self, itinerary, asset_dir, output_path):
        cropped_dir = os.path.join(asset_dir, "smart_cropped")
        if not os.path.exists(cropped_dir): os.makedirs(cropped_dir)
        
        for item in itinerary:
            self._create_slide(item['title'], item['schedule'], item['map_url'], item['photo'], asset_dir, cropped_dir)
        self.prs.save(output_path)
        return True

if __name__ == "__main__":
    asset_dir = r"C:\Users\moos\.gemini\antigravity\brain\2fb7398c-27ff-4ae5-b026-1f183e98301a"
    itinerary = [
        {"title": "05/17 파리 ➔ 니스 TGV 경로", "schedule": "09:09 파리 리옹역 출발 | 14:55 니스역 도착 | TGV 6173편", "map_url": "https://www.google.com/maps/dir/Paris+Gare+de+Lyon/Nice-Ville", "photo": "road_map_day_01_*.png"},
        {"title": "05/21 마르세유 ➔ 카시스 경로", "schedule": "깔랑끄 보트 투어 | 해안도로 드라이브 | 안은희 님 주운전자", "map_url": "https://www.google.com/maps/dir/Marseille/Cassis", "photo": "road_map_day_02_*.png"},
        {"title": "05/24 아를 & 빛의 채석장", "schedule": "레보 드 프로방스 관람 | 랑글루아 다리 | 고흐의 자취", "map_url": "https://www.google.com/maps/dir/Les+Baux-de-Provence/Arles", "photo": "road_map_day_03_*.png"},
        {"title": "05/29 피레네 다르투스테 산악열차", "schedule": "다르투스테 열차 예약 필수 | 피레네 산맥 파노라마", "map_url": "https://www.google.com/maps/search/Petit+Train+d'Artouste", "photo": "road_map_day_04_*.png"},
        {"title": "06/01 보르도 렌트카 반납", "schedule": "샤를라 카네다 점심 | 19:00 보르도 반납 | 보르도 야경", "map_url": "https://www.google.com/maps/dir/Sarlat-la-Can%C3%A9da/Bordeaux", "photo": "road_map_day_05_*.png"}
    ]
    
    out_file = r"C:\Users\moos\Desktop\AHN_EUNHEE_SMART_MAP_MASTERPIECE.pptx"
    engine = SmartFitMapEngine()
    if engine.produce(itinerary, asset_dir, out_file):
        print(f"SMART_FIT_SUCCESS: {os.path.getsize(out_file)} bytes")
