import os
from pptx import Presentation
import json
import sys

sys.stdout.reconfigure(encoding='utf-8')

def extract_styles(pptx_path):
    print(f"Scanning: {pptx_path}")
    prs = Presentation(pptx_path)
    
    layouts = []
    for layout in prs.slide_master.slide_layouts:
        l_data = {"name": layout.name, "placeholders": []}
        for ph in layout.placeholders:
            l_data["placeholders"].append({
                "type": str(ph.placeholder_format.type),
                "idx": ph.placeholder_format.idx,
                "left": ph.left, "top": ph.top, "width": ph.width, "height": ph.height
            })
        layouts.append(l_data)
        
    sample_shapes = []
    if len(prs.slides) > 1:
        slide = prs.slides[1]
        for shape in slide.shapes:
            s_data = {
                "name": shape.name,
                "type": str(shape.shape_type),
                "left": shape.left, "top": shape.top, "width": shape.width, "height": shape.height
            }
            if shape.has_text_frame:
                try:
                    p = shape.text_frame.paragraphs[0]
                    r = p.runs[0]
                    s_data["font"] = {
                        "name": r.font.name,
                        "size": r.font.size.pt if r.font.size else None,
                        "color": str(r.font.color.rgb) if hasattr(r.font.color, 'rgb') else None
                    }
                except: pass
            sample_shapes.append(s_data)
            
    return {"layouts": layouts, "sample_shapes": sample_shapes}

if __name__ == "__main__":
    root = r"d:\아름다운 여행 4.25"
    # Get directories and sort them by name
    dirs = [d for d in os.listdir(root) if os.path.isdir(os.path.join(root, d))]
    
    # Try to find by name first, then by common pattern
    target_dir = None
    for d in dirs:
        if any(c in d for c in ["ߺ", "Gyeom", "견본"]):
            target_dir = os.path.join(root, d)
            break
            
    if not target_dir:
        # Fallback: find the directory with .pptx files
        for d in dirs:
            full_path = os.path.join(root, d)
            if any(f.endswith(".pptx") for f in os.listdir(full_path)):
                target_dir = full_path
                break
            
    if not target_dir:
        print("Template directory not found.")
        sys.exit(1)
        
    pptx_files = [f for f in os.listdir(target_dir) if f.endswith(".pptx")]
    if not pptx_files:
        print(f"No PPTX files found in {target_dir}")
        sys.exit(1)
        
    master_pptx = os.path.join(target_dir, pptx_files[0])
    results = extract_styles(master_pptx)
    
    with open("master_design_specs.json", "w", encoding="utf-8") as f:
        json.dump(results, f, ensure_ascii=False, indent=4)
    print("Master design specifications extracted successfully.")
