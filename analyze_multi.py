import os
from pptx import Presentation
import json

def analyze_multiple_slides(pptx_path, indices=[5, 6, 7]):
    prs = Presentation(pptx_path)
    all_structures = {}
    for idx in indices:
        if idx >= len(prs.slides): continue
        slide = prs.slides[idx]
        structure = []
        for shape in slide.shapes:
            info = {"name": shape.name, "type": str(shape.shape_type), "text": shape.text if shape.has_text_frame else ""}
            structure.append(info)
        all_structures[f"slide_{idx}"] = structure
    
    with open("multi_slide_structure.json", "w", encoding="utf-8") as f:
        json.dump(all_structures, f, indent=4, ensure_ascii=False)

if __name__ == "__main__":
    tpl_path = r"d:\아름다운 여행 4.25\Atelier\templates\260516~0526 안은희 고객님 외 3인 고품격 프랑스 여행 9박 11일(신혜경,이두성 고객님 일정).pptx"
    analyze_multiple_slides(tpl_path)
