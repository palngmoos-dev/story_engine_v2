import os
from pptx import Presentation
import json
import sys

sys.stdout.reconfigure(encoding='utf-8')

def extract_detailed_style(pptx_path):
    prs = Presentation(pptx_path)
    style_data = {
        "file_name": os.path.basename(pptx_path),
        "layouts": [],
        "slides": []
    }

    # Analyze Master Layouts
    for i, layout in enumerate(prs.slide_master.slide_layouts):
        layout_info = {
            "name": layout.name,
            "placeholders": []
        }
        for ph in layout.placeholders:
            layout_info["placeholders"].append({
                "type": str(ph.placeholder_format.type),
                "idx": ph.placeholder_format.idx,
                "left": ph.left,
                "top": ph.top,
                "width": ph.width,
                "height": ph.height
            })
        style_data["layouts"].append(layout_info)

    # Analyze individual slides (first 5 slides for deep sampling)
    for i, slide in enumerate(prs.slides[:5]):
        slide_info = {
            "slide_index": i + 1,
            "layout_name": slide.slide_layout.name,
            "shapes": []
        }
        for shape in slide.shapes:
            shape_info = {
                "name": shape.name,
                "type": str(shape.shape_type),
                "left": shape.left,
                "top": shape.top,
                "width": shape.width,
                "height": shape.height
            }
            if shape.has_text_frame:
                para = shape.text_frame.paragraphs[0] if shape.text_frame.paragraphs else None
                if para:
                    run = para.runs[0] if para.runs else None
                    if run:
                        shape_info["font"] = {
                            "name": run.font.name,
                            "size": run.font.size.pt if run.font.size else None,
                            "bold": run.font.bold,
                            "color": str(run.font.color.rgb) if hasattr(run.font.color, 'rgb') else None
                        }
            
            # Check for arrows/lines (connectors)
            if shape.shape_type == 1: # Line/Connector
                shape_info["line_style"] = {
                    "width": shape.line.width.pt if shape.line.width else None,
                    "color": str(shape.line.color.rgb) if hasattr(shape.line.color, 'rgb') else None
                }
            
            slide_info["shapes"].append(shape_info)
        style_data["slides"].append(slide_info)

    return style_data

if __name__ == "__main__":
    # Get all template files
    template_dir = r"d:\아름다운 여행 4.25\ߺ 4.25"
    if not os.path.exists(template_dir):
        # Try finding by index if path fails due to encoding
        print(f"Directory not found: {template_dir}")
        sys.exit(1)

    pptx_files = [f for f in os.listdir(template_dir) if f.endswith(".pptx")]
    all_analysis = []

    for pptx in pptx_files:
        path = os.path.join(template_dir, pptx)
        print(f"Analyzing {pptx}...")
        try:
            analysis = extract_detailed_style(path)
            all_analysis.append(analysis)
        except Exception as e:
            print(f"Error analyzing {pptx}: {e}")

    with open("master_style_analysis.json", "w", encoding="utf-8") as f:
        json.dump(all_analysis, f, ensure_ascii=False, indent=4)
    print("Master style analysis complete. Saved to master_style_analysis.json")
